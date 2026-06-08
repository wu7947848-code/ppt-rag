"""
pptx_parser.py — python-pptx + lxml chart.xml + openpyxl 嵌入xlsx 解析
"""
from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import List, Optional
from lxml import etree

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
import openpyxl

from backend.index.schema import Page, ChartData

logger = logging.getLogger(__name__)

# XML 命名空间
_NS = {
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}

_CHART_TYPE_MAP = {
    "BAR_CLUSTERED": "bar", "BAR_STACKED": "bar", "BAR_STACKED_100": "bar",
    "COLUMN_CLUSTERED": "bar", "COLUMN_STACKED": "bar",
    "LINE": "line", "LINE_MARKERS": "line", "LINE_STACKED": "line",
    "PIE": "pie", "PIE_EXPLODED": "pie", "DOUGHNUT": "pie",
    "AREA": "area", "AREA_STACKED": "area",
    "XY_SCATTER": "scatter", "XY_SCATTER_LINES": "scatter",
}


def _chart_type_str(chart) -> str:
    try:
        ct = str(chart.chart_type).split(".")[-1]
        return _CHART_TYPE_MAP.get(ct, ct.lower())
    except Exception:
        return "unknown"


def _parse_chart_xml(chart_part) -> tuple[List[str], List[dict]]:
    """从 chart.xml 裸解析 categories 和 series 数据"""
    try:
        xml_bytes = chart_part.blob
        root = etree.fromstring(xml_bytes)
        series_nodes = root.findall(".//c:ser", _NS)

        categories: List[str] = []
        series_list: List[dict] = []

        for ser_node in series_nodes:
            # series name
            ser_name_node = ser_node.find(".//c:tx//c:v", _NS)
            ser_name = ser_name_node.text if ser_name_node is not None else ""

            # categories (从第一个 series 读)
            if not categories:
                cat_nodes = ser_node.findall(".//c:cat//c:v", _NS)
                if not cat_nodes:
                    # numRef categories
                    cat_nodes = ser_node.findall(".//c:cat//c:numRef//c:v", _NS)
                categories = [n.text or "" for n in cat_nodes]

            # values
            val_nodes = ser_node.findall(".//c:val//c:v", _NS)
            if not val_nodes:
                val_nodes = ser_node.findall(".//c:yVal//c:v", _NS)
            values = []
            for vn in val_nodes:
                try:
                    values.append(float(vn.text) if vn.text else None)
                except (ValueError, TypeError):
                    values.append(None)

            series_list.append({"name": ser_name, "values": values})

        return categories, series_list
    except Exception as e:
        logger.debug(f"chart xml parse error: {e}")
        return [], []


def _chart_to_nl(chart_type: str, title: str, categories: List[str],
                 series_list: List[dict]) -> str:
    parts = [f"【{chart_type}图】{title}："]
    for ser in series_list:
        ser_name = ser.get("name", "")
        vals = ser.get("values", [])
        pairs = []
        for i, cat in enumerate(categories[:10]):
            val = vals[i] if i < len(vals) else None
            val_str = str(val) if val is not None else "null"
            pairs.append(f"{cat}={val_str}")
        parts.append(f"{ser_name} " + "，".join(pairs))
    return " ".join(parts)


def _table_to_markdown(table) -> str:
    rows = []
    for i, row in enumerate(table.rows):
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
        if i == 0:
            rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n".join(rows)


def _get_slide_title(slide) -> str:
    """提取幻灯片标题，多级回退"""
    # 1. 标准 title placeholder
    if slide.shapes.title and slide.shapes.title.text.strip():
        return slide.shapes.title.text.strip()

    # 2. placeholder idx=0
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0 and shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                return txt

    # 3. 最大字号的文本框
    best_size = 0
    best_text = ""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    size = run.font.size or 0
                    if size > best_size:
                        best_size = size
                        best_text = run.text.strip()
                except Exception:
                    pass
    return best_text


def _extract_embedded_xlsx(prs) -> List[str]:
    """提取 pptx 包内嵌入的 xlsx 文件，转 markdown 表格"""
    tables_md = []
    try:
        package = prs.part.package
        for part in package.iter_parts():
            if part.partname.endswith(".xlsx"):
                try:
                    wb = openpyxl.load_workbook(io.BytesIO(part.blob), data_only=True)
                    ws = wb.active
                    rows = list(ws.iter_rows(values_only=True, max_row=20))
                    if not rows:
                        continue
                    md_rows = []
                    for i, row in enumerate(rows):
                        cells = [str(c) if c is not None else "" for c in row]
                        md_rows.append("| " + " | ".join(cells) + " |")
                        if i == 0:
                            md_rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                    tables_md.append("\n".join(md_rows))
                except Exception as e:
                    logger.debug(f"embedded xlsx parse error: {e}")
    except Exception as e:
        logger.debug(f"iter_parts error: {e}")
    return tables_md


def parse(pptx_path: Path, doc_id: str, filename: str) -> List[Page]:
    prs = Presentation(str(pptx_path))
    embedded_tables = _extract_embedded_xlsx(prs)
    pages: List[Page] = []

    for slide_idx, slide in enumerate(prs.slides):
        page_num = slide_idx + 1

        title = _get_slide_title(slide)

        text_parts: List[str] = []
        tables_md: List[str] = []
        chart_data_list: List[ChartData] = []
        has_image = False
        has_chart = False
        has_table = False

        for shape in slide.shapes:
            # ---- 图片 ----
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image = True
                continue

            # ---- 表格 ----
            if shape.has_table:
                has_table = True
                try:
                    tables_md.append(_table_to_markdown(shape.table))
                except Exception as e:
                    logger.debug(f"table parse error slide {page_num}: {e}")
                continue

            # ---- 图表 ----
            if shape.has_chart:
                has_chart = True
                try:
                    chart = shape.chart
                    ct = _chart_type_str(chart)

                    # chart title
                    chart_title = ""
                    try:
                        chart_title = chart.chart_title.text_frame.text.strip()
                    except Exception:
                        pass

                    # parse via lxml for reliability
                    categories, series_list = _parse_chart_xml(chart.part)

                    # fallback to python-pptx API
                    if not series_list:
                        for ser in chart.series:
                            try:
                                vals = [v for v in ser.values]
                            except Exception:
                                vals = []
                            series_list.append({"name": getattr(ser, "name", ""), "values": vals})

                    nl = _chart_to_nl(ct, chart_title, categories, series_list)
                    chart_data_list.append(ChartData(
                        chart_type=ct,
                        title=chart_title,
                        categories=categories,
                        series=series_list,
                        nl_description=nl,
                    ))
                except Exception as e:
                    logger.warning(f"chart parse error slide {page_num}: {e}")
                continue

            # ---- 文本框 ----
            if shape.has_text_frame:
                para_texts = []
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        para_texts.append(para_text)
                if para_texts:
                    text_parts.append("\n".join(para_texts))

        text = "\n".join(text_parts).strip()
        chart_nl = " ".join(cd.nl_description for cd in chart_data_list)
        word_count = len(text.split())

        is_divider = (
            word_count < 15
            and not has_image
            and not tables_md
            and not chart_data_list
        )

        pages.append(Page(
            doc_id=doc_id,
            filename=filename,
            page_num=page_num,
            title=title,
            text=text,
            tables_md=tables_md,
            chart_data=chart_data_list,
            chart_nl=chart_nl,
            image_captions=[],
            ocr_texts=[],
            thumbnail_path="",
            is_divider=is_divider,
            has_chart=has_chart,
            has_table=has_table,
            has_image=has_image,
            word_count=word_count,
        ))

    return pages
